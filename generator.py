"""
IB Presentation Generator - Gera apresentações PowerPoint com formatação
de investment banking a partir de dados estruturados em JSON.

Estilo baseado nos templates AGF Advisory / TGA:
- Fonte Inter (Regular, SemiBold, Bold)
- Paleta de cinzas e azul escuro sutil
- Layout limpo e minimalista
"""

import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData


class IBPresentationGenerator:
    """Gera apresentações .pptx com estilo investment banking."""

    # -- Paleta extraída dos templates --------------------------------- #
    IB_BLUE = RGBColor(0x0F, 0x4D, 0x88)
    DARK_TEXT = RGBColor(0x26, 0x26, 0x26)
    TITLE_GRAY = RGBColor(0x7F, 0x7F, 0x7F)
    BODY_GRAY = RGBColor(0x7F, 0x7F, 0x7F)
    LIGHT_TEXT = RGBColor(0xF2, 0xF2, 0xF2)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
    SHAPE_DARK = RGBColor(0x3C, 0x3C, 0x3C)
    SHAPE_MID = RGBColor(0x58, 0x59, 0x5C)
    BORDER_GRAY = RGBColor(0xD9, 0xDA, 0xDD)
    ACCENT_GRAY = RGBColor(0xB3, 0xB5, 0xBB)
    RED_ACCENT = RGBColor(0x98, 0x00, 0x00)
    GREEN_ACCENT = RGBColor(0x1E, 0x82, 0x32)

    CHART_COLORS = [
        RGBColor(0x0F, 0x4D, 0x88),
        RGBColor(0x3C, 0x3C, 0x3C),
        RGBColor(0xB3, 0xB5, 0xBB),
        RGBColor(0x58, 0x59, 0x5C),
        RGBColor(0xD9, 0xDA, 0xDD),
        RGBColor(0x26, 0x26, 0x26),
    ]

    # -- Fontes -------------------------------------------------------- #
    FONT_TITLE = "Inter"
    FONT_BODY = "Inter"

    # -- Dimensões (16:9 padrão PowerPoint) ---------------------------- #
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    MARGIN_L = Inches(0.6)
    MARGIN_R = Inches(0.6)
    CONTENT_W = Inches(12.133)
    TITLE_TOP = Inches(0.35)
    TITLE_H = Inches(0.55)
    SEPARATOR_TOP = Inches(0.95)
    CONTENT_TOP = Inches(1.2)
    CONTENT_H = Inches(5.6)
    FOOTER_TOP = Inches(7.1)

    def __init__(self, template_path: str | None = None):
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            self.prs.slide_width = self.SLIDE_W
            self.prs.slide_height = self.SLIDE_H
        self._page_num = 0

    def generate(self, data: dict) -> bytes:
        slides = data.get("slides", [])
        if not slides:
            raise ValueError("Nenhum slide fornecido no JSON")

        for slide_data in slides:
            layout_name = slide_data.get("layout", "content")
            handler = getattr(self, f"_add_{layout_name}_slide", None)
            if handler is None:
                handler = self._add_content_slide
            handler(slide_data)

        buf = BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    # ================================================================== #
    #  SLIDE TYPES                                                        #
    # ================================================================== #

    def _add_cover_slide(self, data: dict):
        slide = self._new_blank_slide()
        self._set_solid_bg(slide, self.SHAPE_DARK)

        title = data.get("title", "")
        self._add_textbox(
            slide, title,
            self.MARGIN_L, Inches(2.4), self.CONTENT_W, Inches(1.6),
            font_name=self.FONT_TITLE, font_size=Pt(40),
            font_color=self.WHITE, bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        subtitle = data.get("subtitle", "")
        if subtitle:
            self._add_textbox(
                slide, subtitle,
                self.MARGIN_L, Inches(4.2), self.CONTENT_W, Inches(0.9),
                font_name=self.FONT_BODY, font_size=Pt(18),
                font_color=self.ACCENT_GRAY,
                alignment=PP_ALIGN.LEFT,
            )

        date_text = data.get("date", "")
        if date_text:
            self._add_textbox(
                slide, date_text,
                self.MARGIN_L, Inches(6.5), self.CONTENT_W, Inches(0.45),
                font_name=self.FONT_BODY, font_size=Pt(14),
                font_color=self.TITLE_GRAY,
                alignment=PP_ALIGN.LEFT,
            )

    def _add_content_slide(self, data: dict):
        slide = self._new_blank_slide()
        self._set_solid_bg(slide, self.WHITE)
        self._add_title_area(slide, data.get("title", ""))
        self._add_footer(slide)

        bullets = data.get("bullets", [])
        if not bullets:
            return

        top = self.CONTENT_TOP + Inches(0.15)
        for bullet in bullets:
            tb = slide.shapes.add_textbox(
                self.MARGIN_L + Inches(0.3), top,
                self.CONTENT_W - Inches(0.6), Inches(0.65),
            )
            tf = tb.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.space_before = Pt(4)
            p.space_after = Pt(4)

            run_bullet = p.add_run()
            run_bullet.text = "\u2022  "
            run_bullet.font.size = Pt(14)
            run_bullet.font.name = self.FONT_BODY
            run_bullet.font.color.rgb = self.IB_BLUE

            run_text = p.add_run()
            run_text.text = str(bullet)
            run_text.font.size = Pt(14)
            run_text.font.name = self.FONT_BODY
            run_text.font.color.rgb = self.DARK_TEXT

            top += Inches(0.7)

    def _add_two_columns_slide(self, data: dict):
        slide = self._new_blank_slide()
        self._set_solid_bg(slide, self.WHITE)
        self._add_title_area(slide, data.get("title", ""))
        self._add_footer(slide)

        col_w = Inches(5.6)
        gap = Inches(0.9)

        for i, col_key in enumerate(["left_column", "right_column"]):
            col_data = data.get(col_key, {})
            if not col_data:
                col_data = {}
            left = self.MARGIN_L + i * (col_w + gap)

            sub = col_data.get("subtitle", "")
            if sub:
                self._add_textbox(
                    slide, sub,
                    left, self.CONTENT_TOP, col_w, Inches(0.42),
                    font_name=self.FONT_TITLE, font_size=Pt(16),
                    font_color=self.IB_BLUE, bold=True,
                )
                self._add_shape(
                    slide, MSO_SHAPE.RECTANGLE,
                    left, self.CONTENT_TOP + Inches(0.48), col_w, Inches(0.015),
                    fill_color=self.BORDER_GRAY,
                )

            bullets = col_data.get("bullets", [])
            top = self.CONTENT_TOP + Inches(0.65)
            for bullet in bullets:
                tb = slide.shapes.add_textbox(
                    left + Inches(0.15), top,
                    col_w - Inches(0.3), Inches(0.55),
                )
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.space_before = Pt(3)

                run_b = p.add_run()
                run_b.text = "\u2022  "
                run_b.font.size = Pt(12)
                run_b.font.name = self.FONT_BODY
                run_b.font.color.rgb = self.IB_BLUE

                run_t = p.add_run()
                run_t.text = str(bullet)
                run_t.font.size = Pt(13)
                run_t.font.name = self.FONT_BODY
                run_t.font.color.rgb = self.DARK_TEXT
                top += Inches(0.58)

        divider_x = self.MARGIN_L + col_w + gap / 2
        self._add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            divider_x - Inches(0.008), self.CONTENT_TOP,
            Inches(0.016), Inches(5.2),
            fill_color=self.BORDER_GRAY,
        )

    def _add_table_slide(self, data: dict):
        slide = self._new_blank_slide()
        self._set_solid_bg(slide, self.WHITE)
        self._add_title_area(slide, data.get("title", ""))
        self._add_footer(slide)

        tbl_data = data.get("table", {})
        headers = tbl_data.get("headers", [])
        rows = tbl_data.get("rows", [])
        if not headers:
            return

        n_rows = len(rows) + 1
        n_cols = len(headers)

        tbl_height = min(Inches(0.42) * n_rows, Inches(5.2))
        shape = slide.shapes.add_table(
            n_rows, n_cols,
            self.MARGIN_L, self.CONTENT_TOP + Inches(0.15),
            self.CONTENT_W, tbl_height,
        )
        table = shape.table

        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = str(header)
            self._style_cell(
                cell, font_size=Pt(12), bold=True,
                font_color=self.WHITE, fill_color=self.IB_BLUE,
            )

        for ri, row in enumerate(rows):
            bg = self.BG_LIGHT if ri % 2 == 0 else self.WHITE
            for ci, val in enumerate(row):
                cell = table.cell(ri + 1, ci)
                cell.text = str(val)
                fc = self.DARK_TEXT
                if isinstance(val, str):
                    if val.startswith("+"):
                        fc = self.GREEN_ACCENT
                    elif val.startswith("-") and "%" in val:
                        fc = self.RED_ACCENT
                self._style_cell(
                    cell, font_size=Pt(11), font_color=fc, fill_color=bg,
                )

    def _add_chart_slide(self, data: dict):
        slide = self._new_blank_slide()
        self._set_solid_bg(slide, self.WHITE)
        self._add_title_area(slide, data.get("title", ""))
        self._add_footer(slide)

        chart_type_str = data.get("chart_type", "bar")
        chart_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
        }
        xl_type = chart_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        cd = data.get("chart_data", {})
        categories = cd.get("categories", [])
        series_list = cd.get("series", [])
        if not categories or not series_list:
            return

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for s in series_list:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

        chart_left = self.MARGIN_L + Inches(0.5)
        chart_top = self.CONTENT_TOP + Inches(0.2)
        chart_w = self.CONTENT_W - Inches(1.0)
        chart_h = Inches(4.8)

        chart_shape = slide.shapes.add_chart(
            xl_type, chart_left, chart_top, chart_w, chart_h, chart_data,
        )
        chart = chart_shape.chart

        for i, s in enumerate(chart.series):
            if i < len(self.CHART_COLORS):
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = self.CHART_COLORS[i]

        chart.has_legend = len(series_list) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(11)
            chart.legend.font.name = self.FONT_BODY

        if xl_type != XL_CHART_TYPE.PIE:
            cat_axis = chart.category_axis
            cat_axis.tick_labels.font.size = Pt(11)
            cat_axis.tick_labels.font.name = self.FONT_BODY
            cat_axis.tick_labels.font.color.rgb = self.DARK_TEXT
            cat_axis.has_major_gridlines = False

            val_axis = chart.value_axis
            val_axis.tick_labels.font.size = Pt(10)
            val_axis.tick_labels.font.name = self.FONT_BODY
            val_axis.tick_labels.font.color.rgb = self.TITLE_GRAY
            val_axis.major_gridlines.format.line.color.rgb = self.BORDER_GRAY

    def _add_key_metrics_slide(self, data: dict):
        slide = self._new_blank_slide()
        self._set_solid_bg(slide, self.WHITE)
        self._add_title_area(slide, data.get("title", ""))
        self._add_footer(slide)

        metrics = data.get("metrics", [])
        if not metrics:
            return

        n = len(metrics)
        max_per_row = min(n, 4)
        gap = Inches(0.35)
        total_gap = gap * (max_per_row - 1)
        card_w = (self.CONTENT_W - total_gap) / max_per_row
        card_h = Inches(2.2)
        start_y = self.CONTENT_TOP + Inches(0.6)

        for i, metric in enumerate(metrics):
            row = i // max_per_row
            col = i % max_per_row
            left = self.MARGIN_L + col * (card_w + gap)
            top = start_y + row * (card_h + gap)

            self._add_shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, card_w, card_h,
                fill_color=self.BG_LIGHT,
                line_color=self.BORDER_GRAY,
            )

            self._add_shape(
                slide, MSO_SHAPE.RECTANGLE,
                left, top, card_w, Inches(0.04),
                fill_color=self.IB_BLUE,
            )

            label = metric.get("label", "")
            self._add_textbox(
                slide, label,
                left + Inches(0.2), top + Inches(0.25),
                card_w - Inches(0.4), Inches(0.35),
                font_name=self.FONT_BODY, font_size=Pt(12),
                font_color=self.TITLE_GRAY,
                bold=False, alignment=PP_ALIGN.CENTER,
            )

            value = metric.get("value", "")
            self._add_textbox(
                slide, value,
                left + Inches(0.2), top + Inches(0.7),
                card_w - Inches(0.4), Inches(0.65),
                font_name=self.FONT_TITLE, font_size=Pt(26),
                font_color=self.DARK_TEXT, bold=True,
                alignment=PP_ALIGN.CENTER,
            )

            variation = metric.get("variation", "")
            if variation:
                v_color = (
                    self.GREEN_ACCENT if "+" in variation
                    else self.RED_ACCENT if "-" in variation
                    else self.TITLE_GRAY
                )
                self._add_textbox(
                    slide, variation,
                    left + Inches(0.2), top + Inches(1.5),
                    card_w - Inches(0.4), Inches(0.35),
                    font_name=self.FONT_BODY, font_size=Pt(13),
                    font_color=v_color, bold=True,
                    alignment=PP_ALIGN.CENTER,
                )

    def _add_closing_slide(self, data: dict):
        slide = self._new_blank_slide()
        self._set_solid_bg(slide, self.SHAPE_DARK)

        title = data.get("title", "Obrigado")
        self._add_textbox(
            slide, title,
            self.MARGIN_L, Inches(2.8), self.CONTENT_W, Inches(1.2),
            font_name=self.FONT_TITLE, font_size=Pt(42),
            font_color=self.WHITE, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        subtitle = data.get("subtitle", "")
        if subtitle:
            self._add_textbox(
                slide, subtitle,
                self.MARGIN_L, Inches(4.2), self.CONTENT_W, Inches(0.7),
                font_name=self.FONT_BODY, font_size=Pt(16),
                font_color=self.ACCENT_GRAY,
                alignment=PP_ALIGN.CENTER,
            )

    # ================================================================== #
    #  HELPERS                                                            #
    # ================================================================== #

    def _new_blank_slide(self):
        layout = self.prs.slide_layouts[6]  # blank layout
        slide = self.prs.slides.add_slide(layout)
        self._page_num += 1
        return slide

    def _set_solid_bg(self, slide, color: RGBColor):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_title_area(self, slide, title: str):
        """Título cinza em fundo branco com linha separadora fina."""
        self._add_textbox(
            slide, title,
            self.MARGIN_L, self.TITLE_TOP, self.CONTENT_W, self.TITLE_H,
            font_name=self.FONT_TITLE, font_size=Pt(24),
            font_color=self.TITLE_GRAY, bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        self._add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            self.MARGIN_L, self.SEPARATOR_TOP,
            self.CONTENT_W, Inches(0.015),
            fill_color=self.BORDER_GRAY,
        )

    def _add_footer(self, slide):
        self._add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), self.FOOTER_TOP, self.SLIDE_W, Inches(0.01),
            fill_color=self.BORDER_GRAY,
        )

        self._add_textbox(
            slide, str(self._page_num),
            self.SLIDE_W - Inches(1.0), self.FOOTER_TOP + Inches(0.04),
            Inches(0.5), Inches(0.3),
            font_name=self.FONT_BODY, font_size=Pt(10),
            font_color=self.TITLE_GRAY,
            alignment=PP_ALIGN.RIGHT,
        )

    def _add_textbox(self, slide, text, left, top, width, height,
                     font_name=None, font_size=Pt(14), font_color=None,
                     bold=False, alignment=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        p.font.size = font_size
        p.font.name = font_name or self.FONT_BODY
        p.font.bold = bold
        if font_color:
            p.font.color.rgb = font_color
        return tb

    def _add_shape(self, slide, shape_type, left, top, width, height,
                   fill_color=None, line_color=None):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def _style_cell(self, cell, font_size=Pt(12), bold=False,
                    font_color=None, fill_color=None):
        if fill_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.name = self.FONT_BODY
            paragraph.font.bold = bold
            if font_color:
                paragraph.font.color.rgb = font_color
